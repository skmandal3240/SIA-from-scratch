"""SIA AI investor pitch deck — 10 slides. Run: .venv/bin/python build_pitch_deck.py"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parent.parent.parent
OUT = ROOT / "outputs" / "goc"
OUT.mkdir(parents=True, exist_ok=True)

NAVY = RGBColor(0x0F, 0x1B, 0x2D)
GOLD = RGBColor(0xD9, 0xA4, 0x41)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x55, 0x5F, 0x6E)
GREEN = RGBColor(0x2E, 0x7D, 0x32)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def add_text(slide, text, left, top, width, height, size=20, bold=False, color=NAVY):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def title_slide(title, subtitle):
    s = prs.slides.add_slide(blank)
    add_text(s, title, 0.8, 2.2, 11.7, 1.2, size=44, bold=True)
    add_text(s, subtitle, 0.8, 3.6, 11.7, 0.8, size=20, color=GREY)


def bullet_slide(title, bullets, accent=NAVY):
    s = prs.slides.add_slide(blank)
    add_text(s, title, 0.8, 0.5, 11.7, 1.0, size=32, bold=True, color=accent)
    y = 1.6
    for b in bullets:
        add_text(s, f"•  {b}", 0.9, y, 11.5, 0.7, size=18)
        y += 0.72


def split_slide(title, left_title, left_items, right_title, right_items):
    s = prs.slides.add_slide(blank)
    add_text(s, title, 0.8, 0.5, 11.7, 1.0, size=30, bold=True)
    add_text(s, left_title, 0.8, 1.7, 5.5, 0.6, size=20, bold=True, color=GOLD)
    y = 2.4
    for b in left_items:
        add_text(s, f"•  {b}", 0.9, y, 5.6, 0.6, size=16)
        y += 0.58
    add_text(s, right_title, 7.0, 1.7, 5.5, 0.6, size=20, bold=True, color=GOLD)
    y = 2.4
    for b in right_items:
        add_text(s, f"•  {b}", 7.1, y, 5.6, 0.6, size=16)
        y += 0.58


title_slide("SIA — Private On-Device AI for India", "India-first, from-scratch multimodal AI | Saurabh Mandal | Aug 2026")

bullet_slide("The Problem", [
    "900M+ Indians lack reliable high-speed internet — cloud AI is unusable offline",
    "Privacy: users' data leaves their devices to foreign servers (DPDP 2023 pressure)",
    "No credible Indian-language, on-device AI assistant exists today",
])

bullet_slide("The Solution: SIA", [
    "A private AI companion that runs entirely on your device — no cloud, no data leaving",
    "Multimodal by design: text, audio (listen + generate), vision, video, code",
    "From-scratch stack: tokenizer → transformer → tools → training, all original code",
    "Built in India, for Indian languages (Hindi + English first)",
])

split_slide("Why We Win",
    "Tech moat", [
        "From-scratch transformer — no API wrapper fees",
        "On-device privacy is the product, not an add-on",
        "Audio-native: hears music/voice, speaks back",
        "Edge INT4/INT8 quantized models for 8GB devices",
    ],
    "Market gap", [
        "Cloud chatbots can't work offline",
        "India AI market growing >25%/yr",
        "No credible Indian on-device multimodal assistant",
        "B2C freemium + B2B white-label edge AI",
    ])

bullet_slide("Traction — Working Demos", [
    "From-scratch transformer trained on CPU: 4,000 steps, loss 2.78",
    "9-demo gauntlet passing: text, code, vision, audio-listen, audio-gen, image, video, tools, quantize",
    "Audio→language loop live: SIA hears a WAV, describes it, speaks a reply",
    "VAE decoder: text → 256×256 RGB image generation",
    "Tool-use agent loop: calc/now/file tools via [[tool:name(args)]]",
])

split_slide("Business Model",
    "Revenue (Pillar 1)", [
        "Freemium companion app (privacy = premium)",
        "B2B white-label edge AI for Indian OEMs",
        "Reinvest 30–50% margin into models + compute",
    ],
    "Non-dilutive fuel", [
        "TIDE 2.0 (MeitY) — up to ₹50L, dossier ready",
        "SISFS (DST) — up to ₹50L",
        "IndiaAI compute credits (40% off)",
        "Bihar Startup Policy — Patna-based entity",
    ])

bullet_slide("Roadmap (12 months)", [
    "M1 (0–3): P1 LoRA fine-tune ≥95% tool-call accuracy; SIA Edge INT8 demo on 8GB device",
    "M2 (3–6): First revenue app live — vernacular freemium assistant, 1k users",
    "M3 (6–9): B2B pilot with one Indian OEM; SIA Pro multilingual models",
    "M4 (9–12): 10k users; edge ASIC pre-feasibility (DLI/PLI)",
])

bullet_slide("The Ask", [
    "TIDE 2.0: ₹50L — GPU compute, Indian-language data, product dev, team",
    "Compute: ~100–200 GPU-hours/month for LoRA + distillation",
    "Incubator partner: 51 TIDE incubators — Bihar/Bengaluru",
    "Demo video available on request",
])

bullet_slide("Why Now", [
    "IndiaAI Mission: 34,000 GPUs, 20 sovereign models being funded — timing is right",
    "DPDP 2023 makes on-device privacy a regulatory advantage",
    "Jevons paradox: cheaper edge AI → exploding usage",
    "We already have a working from-scratch stack — this is not a deck-only startup",
])

s = prs.slides.add_slide(blank)
add_text(s, "SIA — Your private intelligence, on your device.", 0.8, 3.0, 11.7, 1.0, size=36, bold=True, color=NAVY)
add_text(s, "saurabh@siacompany.in  |  github.com/skmandal3240/SIA-from-scratch", 0.8, 4.3, 11.7, 0.7, size=18, color=GREY)

prs.save(OUT / "SIA_AI_PITCH_DECK.pptx")
print(f"Saved {OUT / 'SIA_AI_PITCH_DECK.pptx'} ({len(prs.slides.slides if hasattr(prs.slides,'slides') else prs.slides._sldIdLst)} slides)")
