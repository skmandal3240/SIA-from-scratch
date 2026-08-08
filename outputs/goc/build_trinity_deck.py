#!/usr/bin/env python3
"""Project TRINITY — updated investor pitch deck (16 slides) with all new details.
Build: python3 outputs/goc/build_trinity_deck.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parent.parent.parent
OUT = ROOT / "outputs" / "goc" / "TRINITY_Pitch_Deck.pptx"
DIAGRAMS = ROOT / "outputs" / "goc" / "diagrams"

NAVY = RGBColor(0x0F, 0x1B, 0x2D)
GOLD = RGBColor(0xD9, 0xA4, 0x41)
GREY = RGBColor(0x55, 0x5F, 0x6E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def add_text(slide, text, left, top, width, height, size=20, bold=False, color=NAVY):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def title_slide(title, subtitle):
    s = prs.slides.add_slide(blank)
    add_text(s, title, 0.8, 2.3, 11.7, 1.2, size=46, bold=True)
    add_text(s, subtitle, 0.8, 3.7, 11.7, 0.8, size=18, color=GREY)


def bullet_slide(title, bullets):
    s = prs.slides.add_slide(blank)
    add_text(s, title, 0.8, 0.5, 11.7, 1.0, size=30, bold=True)
    y = 1.6
    for b in bullets:
        add_text(s, f"•  {b}", 0.9, y, 11.5, 0.7, size=17)
        y += 0.7


def split_slide(title, l_title, l_items, r_title, r_items):
    s = prs.slides.add_slide(blank)
    add_text(s, title, 0.8, 0.5, 11.7, 0.9, size=28, bold=True)
    add_text(s, l_title, 0.8, 1.6, 5.5, 0.6, size=18, bold=True, color=GOLD)
    y = 2.3
    for b in l_items:
        add_text(s, f"•  {b}", 0.9, y, 5.6, 0.6, size=15)
        y += 0.6
    add_text(s, r_title, 7.0, 1.6, 5.5, 0.6, size=18, bold=True, color=GOLD)
    y = 2.3
    for b in r_items:
        add_text(s, f"•  {b}", 7.1, y, 5.6, 0.6, size=15)
        y += 0.6


def pic_slide(title, img, caption):
    s = prs.slides.add_slide(blank)
    add_text(s, title, 0.8, 0.4, 11.7, 0.8, size=28, bold=True)
    p = DIAGRAMS / img
    if p.exists():
        s.shapes.add_picture(str(p), Inches(1.5), Inches(1.3), width=Inches(10.3))
    add_text(s, caption, 0.8, 6.9, 11.7, 0.5, size=12, color=GREY)


image_slide = pic_slide


# 1. Title
title_slide("Project TRINITY",
            "India's Sovereign AI Infrastructure Company | SIA / Mandal Holdings | August 2026")

# 2. The thesis
bullet_slide("The Thesis", [
    "India needs its own complete AI stack — compute, data, models, agents, applications — not dependence on foreign cloud",
    "TRINITY builds all five layers; SIA is the AI operating system delivered on every device (Nano → Cloud)",
    "One profitable app funds all R&D — sell outcomes, not tokens (Pillar 1 rule)",
    "Phase-gated funding: MVP ₹2–10 Cr → Global ₹2,000 Cr+; registered Bengaluru, Karnataka",
])

# 3. Problem
bullet_slide("The Problem", [
    "Limited sovereign AI infrastructure; dependence on foreign foundation models",
    "900M+ Indians lack reliable connectivity — cloud AI is unusable offline",
    "DPDP 2023: data leaving devices is a compliance risk — on-device is the answer",
    "High compute costs, fragmented tooling, no credible Indian full-stack player",
])

# 4. The Solution — SIA
bullet_slide("The Solution — SIA, an AI Operating System", [
    "Family of models on ONE shared stack: tokenizer, architecture, tool-calling, memory format, runtime",
    "SIA Nano 0.5–1B (wearables/IoT) · SIA Edge 2–4B (phones) · SIA Pro 8–14B (workstations) · SIA Cloud 70B+ (data centers)",
    "Private by design: on-device inference keeps data in India",
    "Multimodal: text, audio, vision, code — audio listen/generate live; VAE image decode working",
])

# 5. Pillars
split_slide("India's Five-Layer AI Prize",
    "The stack", [
        "1. Compute — GPU cloud, chips (GPU/TPU/NPU/photonic)",
        "2. Data — Indian languages, synthetic, SIA Memory",
        "3. Models — SIA family + services",
        "4. Agents & OS — AI runtime, SIA Studio/SDK",
        "5. Applications — health/edu/agri/defence/gov",
    ],
    "Why it matters", [
        "Each layer earns its own margin",
        "Sovereign moat — replacing foreign dependency",
        "Jobs, IP, chips, energy stay in India",
        "Government is both customer and mandate",
    ])

# 6. Architecture (diagram 1)
image_slide("The 5-Layer AI Stack", "01_5layer_ai_stack.png",
            "L1 coding agents → L5 chip+energy; outcomes: SIA Everywhere · Root India / Build India")

# 7. Model family (diagram 2)
image_slide("SIA Model Family — One OS, Four Tiers", "02_sia_model_family.png",
            "Nano→Cloud share tokenizer, architecture, tool interface, memory formats, common runtime")

# 8. Memory/agent loop (diagram 3)
image_slide("The Agent & Memory Loop", "03_memory-agents-loop.png",
            "User → Agent → Tools → Core → Model Family → Memory · DPDP privacy boundary")

# 9. Market
bullet_slide("Market — TAM/SAM/SOM", [
    "TAM: India AI spend $12–15B (2026E) → $25–40B (2030E); global AI $800B+ by 2030",
    "SAM: ~$4–6B by 2028 — enterprise AI + on-device Indic + SDK/embedded",
    "SOM: 1% by Y5 ≈ $40–60M ARR (≈₹330–500 Cr)",
    "First-mover: no credible full-stack sovereign Indian AI company today",
])

# 10. Business model + financials
bullet_slide("Business Model & Financial Profile", [
    "Revenue: API · enterprise subscriptions · GPU cloud · licenses · managed services",
    "Indicative build: ₹0 → ~₹11 Cr (Y3) → ~₹68 Cr (Y5) → ~₹460 Cr (Y10)",
    "Gross margin 70–85%; LTV/CAC target >3 (4–6×); enterprise churn <5%",
    "Phase-gated raises: MVP ₹2–10 Cr → Growth ₹25–100 Cr → Scale ₹100–500 Cr → National ₹500–2,000 Cr",
])

# 11. Competitive
bullet_slide("Competitive Landscape — We Own the Gap", [
    "Frontier labs (OpenAI/Anthropic/Google): cloud-first, English-first, no sovereign-India",
    "Meta/OSS: strong weights, foreign lineage",
    "India labs: fragmented, not full-stack",
    "TRINITY: from-scratch IP, on-device privacy, Indic-first, sovereign-government alignment — the only full-stack play",
])

# 12. Why India / why now
bullet_slide("Why India, Why Now", [
    "IndiaAI Mission: 34,000+ GPUs, 367+ datasets, 20 sovereign models already funded [VERIFIED]",
    "DPDP 2023 makes on-device privacy a regulatory advantage",
    "Dholera 28nm fab + DLI/PLI → silicon roadmap is real",
    "1.4B people; 900M+ offline — on-device is the only deliverable path to them",
])

# 13. Roadmap
bullet_slide("Roadmap (2026–2035)", [
    "Phase 1 (26–27): MVP — SIA Edge companion + 5 pilots",
    "Phase 2 (27–29): Foundation models — Nano/Edge trained, Pro LoRA",
    "Phase 3 (29–31): Enterprise platform — agents, SDK, memory",
    "Phase 4 (31–35): National infra — private cluster → hyperscale campus",
    "Phase 5 (35+): Global expansion",
])

# 14. Traction
bullet_slide("Traction — Working Proof", [
    "From-scratch nano transformer trained on CPU — 9/9 demos pass (text, code, vision, audio-listen, audio-gen, image, video, tools, quantize)",
    "Audio→language loop live; VAE decoder gives real 256×256 images",
    "LoRA fine-tune pipeline ready (Gemma 3 1B, Colab T4)",
    "Full investor doc set: master report, A/B/C appendices, TAM, financials, competitive, GTM, hiring, valuation",
])

# 15. The Ask
bullet_slide("The Ask — ₹2–10 Cr MVP", [
    "Use of funds: team ~30%, compute ~35%, data/fine-tuning ~20%, legal/SDK ~15%",
    "Targets: live flagship demo + 3–5 enterprise/government pilots + DPIIT/Startup India registration",
    "Milestones unlock Growth round ₹25–100 Cr for foundation models",
    "Government path: TIDE 2.0, SISFS, IndiaAI, SAMRIDH + Karnataka state schemes",
])

# 16. Closing
title_slide("Join the Sovereign AI Builder",
           "SIA — Your private intelligence, on your device. India is the root.")

prs.save(OUT)
print(f"Saved {OUT} ({len(prs.slides)} slides)")