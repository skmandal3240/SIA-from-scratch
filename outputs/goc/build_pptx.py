"""Build the GoC 5-pillar PPTX deck from the plan doc. Run: .venv/bin/python build_pptx.py"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parent.parent.parent  # repo root
OUT = ROOT / "outputs" / "goc"
OUT.mkdir(parents=True, exist_ok=True)

NAVY = RGBColor(0x0F, 0x1B, 0x2D)
GOLD = RGBColor(0xD9, 0xA4, 0x41)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x55, 0x5F, 0x6E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def add_text(slide, text, left, top, width, height, size=20, bold=False, color=NAVY):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def title_slide(title, subtitle):
    s = prs.slides.add_slide(blank)
    add_text(s, title, 0.8, 2.2, 11.7, 1.2, size=44, bold=True)
    add_text(s, subtitle, 0.8, 3.6, 11.7, 0.8, size=20, color=GREY)


def bullet_slide(title, bullets):
    s = prs.slides.add_slide(blank)
    add_text(s, title, 0.8, 0.5, 11.7, 1.0, size=32, bold=True)
    y = 1.6
    for b in bullets:
        add_text(s, f"•  {b}", 0.9, y, 11.5, 0.7, size=18)
        y += 0.72


title_slide("SIA — AI 5-Pillar Plan", "India-first, self-funded | Saurabh Mandal Holdings | Aug 2026")

bullet_slide("The Thesis", [
    "OpenAI's moat = inference cost (Jevons paradox: cheaper AI -> exploding usage)",
    "SIA answers India-style: cheaper models, Indian silicon, rented Indian compute now, own DC later",
    "One profitable app funds all R&D — sell outcomes, not tokens",
])

bullet_slide("Pillar 5 — Energy", [
    "PFBR Kalpakkam 500 MWe first criticality 6 Apr 2026 (BHAVINI/DAE)",
    "Nuclear Mission: 100 GW by 2047; Bharat Small Reactors (220 MWe) opening to private sector",
    "Near-term: quantised edge models = 10-50x less energy per inference",
    "12-24 mo: green PPA (IEX/SECI), DC near cheap power, PUE < 1.3",
    "2030s: apply as anchor tenant for dedicated nuclear allocation (NPCIL)",
])

bullet_slide("Pillar 4 — Hardware", [
    "Tata-PSMC Dholera 28nm fab: trial production Dec 2026, ~50K wafers/mo by 2028",
    "Design our own RISC-V + NPU edge accelerator (28nm) for cameras/drones/phones",
    "Partners: InCore/C-DAC (RISC-V), Tata Electronics (fab), DLI covers ~50% design cost",
    "Hardware-agnostic software layer: ONNX Runtime / llama.cpp / vLLM",
    "LPU-class SRAM inference adopted now; optical = 2028+ R&D with IITs/IISc",
])

bullet_slide("Pillar 3 — Infrastructure", [
    "IndiaAI Mission: 38,000+ GPUs; startups up to 40% off -> Rs 65-92/hr",
    "Empaneled: Jio, Yotta, E2E, Tata Comms, C-DAC PARAM (compute.indiaai.gov.in)",
    "Run ALL training on rented Indian GPU (portable: K8s/Slurm + containers)",
    "Colocate when GPU spend > Rs 10-15L/month sustained",
    "Own DC trigger: revenue covers 100% capex or grant secured",
])

bullet_slide("Pillar 2 — Models", [
    "SIA-lite (0.5-1B): cameras, drones, IoT | SIA-edge (3-8B): edge boxes | SIA-pro (70B-class): cloud",
    "From-scratch transformer core IN PROGRESS — tokenizer -> attention -> training (this repo)",
    "Multimodal in stages: text -> vision -> audio -> video/code",
    "Mo 3-6: LoRA fine-tune 7-8B open model on Indian data (~Rs 20-60K)",
    "Do NOT pretrain 70B from scratch yet — $10M+ and no moat",
])

bullet_slide("Pillar 1 — Applications (revenue engine)", [
    "Chosen app #1: Vernacular AI support/agent platform (Hindi + regional voice/text)",
    "Every Indian company has a support budget; sell per-resolution; pure SaaS",
    "MVP this quarter on SIA-pro; 1 anchor enterprise; bill via UPI",
    "Mo 6+: app #2 (drone/camera edge analytics) from same model family",
    "Entities: SIA AI Pvt Ltd (models) + Mandal Devices (hw) + Akasha Runtime (agents)",
])

bullet_slide("Execution Order & Cost", [
    "Phase 0 (wk 1-4): IndiaAI credits + grants; finish SIA core; POC to 5 customers",
    "Phase 1 (mo 2-6): app #1 live + revenue; SIA-pro LoRA; edge quantisation",
    "Phase 2 (mo 6-12): app #2; SIA-edge on hardware; RISC-V SoC design start",
    "Phase 3 (yr 2): tape-out edge ASIC (Rs 15-30 Cr shared); colocation",
    "Phase 4 (yr 3-5): own DC near green power; optical R&D; BSR application; Global South",
])

bullet_slide("First 7 Days", [
    "Apply IndiaAI compute portal (need DPIIT/Startup India number)",
    "Apply SAMRIDH / TIDE 2.0 / IndiaAI grants",
    "Ship SIA transformer core smoke test — ON TRACK",
    "List 20 candidate app #1 customers; call 5",
    "Draft intro emails: Tata Electronics + InCore; read PFBR/BSR policy",
])

out_path = OUT / "SIA_5_Pillar_Plan.pptx"
prs.save(str(out_path))
print(f"PPTX saved: {out_path} ({len(prs.slides._sldIdLst)} slides)")
